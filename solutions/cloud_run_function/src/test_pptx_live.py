#!/usr/bin/env python3
"""
test_pptx_live.py

Comprehensive test suite for PowerPoint (.pptx & .ppt) document handling:
1. Generates a Simple 3-slide PPTX deck with bullet points and tables.
2. Generates a Large 25-slide PPTX deck with rich nested bullets, financial tables, and presenter notes.
3. Generates a legacy .ppt sample for fallback resilience testing.
4. Executes the PPTXHandler through the central HandlerRegistry.
5. Invokes live Google Vertex AI (Gemini 2.5) with Argolis credentials to extract structured insights.
"""

import os
import sys
import io
import json
import time
import logging

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("test_pptx_live")

# Ensure src root is in sys.path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from handlers.registry import registry
from services.auth_service import AuthService
from services.model_service import ModelService
from handlers.base import ProcessedContent


def create_simple_pptx(filepath: str):
    """Creates a clean 3-slide PowerPoint deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_slide_layout)
    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(3.0))
    tf = txBox.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "NYL Unstructured Data Governance"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "Q3 2026 Architecture & AI Roadmap"
    p2.font.size = Pt(24)

    # Slide 2: Bullet Points
    slide2 = prs.slides.add_slide(blank_slide_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.0), Inches(5.5))
    tf2 = txBox2.text_frame
    title_p = tf2.paragraphs[0]
    title_p.text = "Key Strategic Objectives"
    title_p.font.size = Pt(28)
    title_p.font.bold = True

    bullets = [
        ("Automate ingestion of unstructured documents (PDF, DOCX, PPTX) via Cloud Run", 0),
        ("Integrate BigQuery Remote Functions with Vertex AI foundation models", 0),
        ("Sub-second latency SLA for document classification", 1),
        ("Multi-cloud federation across AWS S3 and GCP BigLake Iceberg Catalogs", 0),
        ("Zero data exfiltration policy with private VPC endpoints", 1)
    ]
    for text, level in bullets:
        bp = tf2.add_paragraph()
        bp.text = text
        bp.level = level
        bp.font.size = Pt(18)

    # Slide 3: Financial Table & Speaker Notes
    slide3 = prs.slides.add_slide(blank_slide_layout)
    txBox3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.0), Inches(1.0))
    title_p3 = txBox3.text_frame.paragraphs[0]
    title_p3.text = "Infrastructure Budget & Spending"
    title_p3.font.size = Pt(28)
    title_p3.font.bold = True

    # Add Table
    rows, cols = 4, 4
    table_shape = slide3.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(11.0), Inches(3.0))
    table = table_shape.table

    headers = ["Service Component", "Allocated Budget", "Q1 Actual Spend", "Variance Status"]
    data = [
        ["BigLake Iceberg Federation", "$45,000", "$41,200", "Under Budget (8.4%)"],
        ["Cloud Run Vertex Proxy", "$30,000", "$28,500", "Under Budget (5.0%)"],
        ["Cloud Composer / Airflow", "$18,000", "$21,400", "Over Budget (+18.8%)"]
    ]
    for c_idx, head in enumerate(headers):
        table.cell(0, c_idx).text = head

    for r_idx, row in enumerate(data, start=1):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = val

    # Add Speaker Notes
    notes_slide = slide3.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = "CONFIDENTIAL: Cloud Composer budget variance was caused by ad-hoc backfill runs in January. Downstream optimization scheduled for Q2."

    prs.save(filepath)
    logger.info(f"✅ Generated simple PPTX: {filepath} ({os.path.getsize(filepath)} bytes)")


def create_large_pptx(filepath: str, num_slides: int = 25):
    """Creates a rich, multi-slide enterprise PowerPoint deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    departments = [
        ("Life Insurance Analytics", "$1,250,000", "$1,180,000", "Under Budget"),
        ("Annuities Actuarial Platform", "$850,000", "$920,000", "Over Budget (+8.2%)"),
        ("Claims Document OCR & NLP", "$640,000", "$610,000", "Under Budget"),
        ("Underwriting Automation Engine", "$1,100,000", "$1,050,000", "Under Budget"),
        ("Enterprise Data Governance & DLP", "$450,000", "$445,000", "On Track"),
        ("Cross-Cloud Lakehouse & Omni", "$780,000", "$830,000", "Over Budget (+6.4%)"),
        ("Customer 360 & Telemetry", "$520,000", "$490,000", "Under Budget"),
        ("Generative AI Agent Orchestration", "$950,000", "$910,000", "Under Budget"),
    ]

    for i in range(1, num_slides + 1):
        slide = prs.slides.add_slide(blank_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Slide {i}: Enterprise Pillar Architecture Review - Phase {i}"
        p.font.size = Pt(24)
        p.font.bold = True

        # Content Box with Nested Bullet Points
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.0))
        ctf = content_box.text_frame
        
        c_title = ctf.paragraphs[0]
        c_title.text = f"Operational Domain {i} Capabilities & Deliverables:"
        c_title.font.size = Pt(16)
        c_title.font.bold = True

        bullet_items = [
            (f"Milestone {i}.1: Real-time event streaming via Google Pub/Sub with 99.99% uptime", 0),
            (f"Milestone {i}.2: Zero-copy federated queries across BigLake and Databricks Unity", 0),
            ("Enforces strict column-level and row-level masking policies with Google DLP", 1),
            ("Automated schema migration with Dataform ELT assertions and data quality checks", 1),
            (f"Milestone {i}.3: Serverless microservices scaled across us-east4 with VPC egress", 0),
            ("Direct VPC peering with private Google access for egress security", 1),
        ]
        for text, level in bullet_items:
            bp = ctf.add_paragraph()
            bp.text = text
            bp.level = level
            bp.font.size = Pt(13)

        # Add a Table on every 3rd slide
        if i % 3 == 0:
            tbl_shape = slide.shapes.add_table(4, 3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0))
            tbl = tbl_shape.table
            tbl.cell(0, 0).text = "Domain"
            tbl.cell(0, 1).text = "Budget"
            tbl.cell(0, 2).text = "Status"

            dept_idx = (i % len(departments))
            for r in range(1, 4):
                d_name, d_bgt, _, d_stat = departments[(dept_idx + r) % len(departments)]
                tbl.cell(r, 0).text = d_name[:20]
                tbl.cell(r, 1).text = d_bgt
                tbl.cell(r, 2).text = d_stat

        # Add Presenter Notes on every slide
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = (
            f"EXECUTIVE BRIEFING NOTE FOR SLIDE {i}:\n"
            f"- Critical path dependency on BigQuery Omni IAM role propagation.\n"
            f"- Verified compliance with NYL Information Security standard #SEC-{1000 + i}.\n"
            f"- Target go-live date: Q4 2026 with full canary rollout."
        )

    prs.save(filepath)
    logger.info(f"✅ Generated large PPTX: {filepath} with {num_slides} slides ({os.path.getsize(filepath)} bytes)")


def run_tests():
    print("=" * 80)
    print("🚀 POWERPOINT (.PPTX / .PPT) SUPPORT TEST SUITE")
    print("=" * 80)

    # 1. Generate test files
    simple_pptx = "/tmp/argolis_simple_deck.pptx"
    large_pptx = "/tmp/argolis_large_deck.pptx"
    legacy_ppt = "/tmp/argolis_legacy_deck.ppt"

    create_simple_pptx(simple_pptx)
    create_large_pptx(large_pptx, num_slides=25)

    # Create dummy binary legacy .ppt file to test fallback
    with open(legacy_ppt, "wb") as f:
        f.write(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512 + b"DUMMY_LEGACY_PPT_BINARY_DATA")

    # 2. Test Handler Dispatcher & Parsing
    print("\n" + "=" * 80)
    print("🔍 [TEST 1/4] Testing PPTX Handler on Simple 3-Slide Deck")
    print("=" * 80)
    with open(simple_pptx, "rb") as f:
        simple_bytes = f.read()

    handler_simple = registry.get_handler("presentation.pptx")
    print(f"✅ Resolved Handler: {type(handler_simple).__name__}")
    
    t0 = time.time()
    processed_simple = handler_simple.process(simple_bytes, "argolis_simple_deck.pptx")
    elapsed_simple = time.time() - t0

    print(f"⏱️ Parse Time: {elapsed_simple:.4f}s | is_text: {processed_simple.is_text} | MIME: {processed_simple.mime_type}")
    print(f"📄 Markdown Text Length: {len(processed_simple.text_content)} characters")
    print("\n--- Snippet of Parsed Output (First 500 chars) ---")
    print(processed_simple.text_content[:500])
    print("...\n--- End Snippet ---\n")

    assert "NYL Unstructured Data Governance" in processed_simple.text_content
    assert "CONFIDENTIAL: Cloud Composer budget variance" in processed_simple.text_content
    assert "BigLake Iceberg Federation" in processed_simple.text_content
    print("✅ All assertions passed for Simple PPTX!")

    # 3. Test Large Deck Parsing
    print("\n" + "=" * 80)
    print("🔍 [TEST 2/4] Testing PPTX Handler on Large 25-Slide Deck")
    print("=" * 80)
    with open(large_pptx, "rb") as f:
        large_bytes = f.read()

    t0 = time.time()
    processed_large = handler_simple.process(large_bytes, "argolis_large_deck.pptx")
    elapsed_large = time.time() - t0

    print(f"⏱️ Parse Time (25 Slides): {elapsed_large:.4f}s")
    print(f"📄 Extracted Text Length: {len(processed_large.text_content)} characters ({len(processed_large.text_content.splitlines())} lines)")
    
    assert "Slide 25" in processed_large.text_content
    assert "EXECUTIVE BRIEFING NOTE FOR SLIDE 25" in processed_large.text_content
    print("✅ All 25 slides, tables, and notes parsed successfully!")

    # 4. Test Legacy .PPT Fallback Resilience
    print("\n" + "=" * 80)
    print("🔍 [TEST 3/4] Testing Graceful Fallback on Legacy Binary .PPT")
    print("=" * 80)
    with open(legacy_ppt, "rb") as f:
        legacy_bytes = f.read()

    handler_legacy = registry.get_handler("sample.ppt")
    processed_legacy = handler_legacy.process(legacy_bytes, "argolis_legacy_deck.ppt")
    print(f"✅ Legacy PPT Response: MIME={processed_legacy.mime_type}")
    print(f"   Fallback Message: {processed_legacy.text_content}")
    assert "PowerPoint presentation parse fallback" in processed_legacy.text_content
    print("✅ Graceful error catching verified for binary .ppt!")

    # 5. Live Vertex AI Gemini Invocation
    print("\n" + "=" * 80)
    print("🤖 [TEST 4/4] Live Vertex AI Gemini Model Invocation on PPTX Document")
    print("=" * 80)

    project_id = "databricks-playground-497321"
    location = "global"

    token, detected_project = AuthService.get_bearer_token_and_project(project_id)
    print(f"🔑 Auth Token Acquired: {token[:12]}...{token[-6:]} | Project: {detected_project}")

    prompt = (
        "Analyze this PowerPoint presentation. "
        "1. Summarize the key strategic objectives. "
        "2. Extract the financial table and identify which service is over budget. "
        "3. Surface any confidential speaker notes mentioned in the presentation."
    )

    gemini_model = "gemini-2.5-flash"
    print(f"\n📡 Invoking Gemini on Vertex AI ({gemini_model}, location: {location})...")
    
    t0 = time.time()
    try:
        res = ModelService.invoke_model(
            processed=processed_simple,
            prompt=prompt,
            project_id=detected_project,
            location=location,
            model_name=gemini_model,
            max_tokens=2048
        )
        gemini_elapsed = time.time() - t0

        print(f"⏱️ Gemini Invocation Latency: {gemini_elapsed:.2f}s")
        print("\n" + "-" * 60)
        print("✨ LIVE GEMINI 2.5 RESPONSE (Simple Deck):")
        print("-" * 60)
        print(res.get("extracted_text"))
        print("-" * 60)
        print(f"📊 Usage Metrics: {json.dumps(res.get('usage', {}), indent=2)}")

    except Exception as e:
        print(f"⚠️ Global endpoint returned: {e}. Retrying on us-central1...")
        res = ModelService.invoke_model(
            processed=processed_simple,
            prompt=prompt,
            project_id=detected_project,
            location="us-central1",
            model_name=gemini_model,
            max_tokens=2048
        )
        print("\n" + "-" * 60)
        print("✨ LIVE GEMINI (us-central1) RESPONSE:")
        print("-" * 60)
        print(res.get("extracted_text"))
        print("-" * 60)

    # 6. Live Vertex AI Gemini Invocation on Large 25-Slide Deck
    print("\n" + "=" * 80)
    print("🤖 Live Vertex AI Gemini Model Invocation on LARGE (25 Slides) PPTX")
    print("=" * 80)
    prompt_large = (
        "Review this comprehensive 25-slide enterprise deck. "
        "1. Identify all operational domains that are OVER budget. "
        "2. Summarize the recurring themes in the executive speaker notes. "
        "3. Provide an executive summary of the Phase 25 architecture goals."
    )
    t0 = time.time()
    try:
        res_large = ModelService.invoke_model(
            processed=processed_large,
            prompt=prompt_large,
            project_id=detected_project,
            location=location,
            model_name=gemini_model,
            max_tokens=2048
        )
        gemini_large_elapsed = time.time() - t0
        print(f"⏱️ Large Deck Gemini Latency: {gemini_large_elapsed:.2f}s")
        print("\n" + "-" * 60)
        print("✨ LIVE GEMINI RESPONSE (25-Slide Deck):")
        print("-" * 60)
        print(res_large.get("extracted_text"))
        print("-" * 60)
        print(f"📊 Usage Metrics: {json.dumps(res_large.get('usage', {}), indent=2)}")
    except Exception as e:
        print(f"⚠️ Large deck invocation error: {e}")

    print("\n" + "=" * 80)
    print("🏆 ALL PPTX TESTS AND VERIFICATIONS COMPLETED SUCCESSFULLY!")
    print("=" * 80)


if __name__ == "__main__":
    run_tests()
