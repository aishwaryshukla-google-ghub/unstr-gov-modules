import io
from .base import BaseDataHandler, ProcessedContent


class PPTXHandler(BaseDataHandler):
    """
    Handler for Microsoft PowerPoint (.pptx) presentations.
    Extracts slide titles, body text, bullet points, tables, and speaker notes
    directly into structured Markdown text for LLM processing.
    """

    def process(self, file_bytes: bytes, filename: str) -> ProcessedContent:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            md_lines = [f"--- START OF POWERPOINT PRESENTATION ({filename}) ---", ""]

            for slide_num, slide in enumerate(prs.slides, start=1):
                md_lines.append(f"## Slide {slide_num}")

                # 1. Extract Slide Title if present
                title_shape = slide.shapes.title
                if title_shape and title_shape.text.strip():
                    md_lines.append(f"**Title**: {title_shape.text.strip()}\n")

                # 2. Extract Shape Text, Bullets, and Tables
                for shape in slide.shapes:
                    # Skip the title shape since it's already rendered above
                    if shape == title_shape:
                        continue

                    # Handle Text Frames (body copy, callouts, text boxes)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if not text:
                                continue
                            indent = "  " * paragraph.level
                            # Prefix with bullet if it has an indentation level or multiple items
                            bullet = "- " if paragraph.level > 0 else ""
                            md_lines.append(f"{indent}{bullet}{text}")

                    # Handle Tables
                    elif shape.has_table:
                        table = shape.table
                        md_lines.append("\n### Table on Slide:")
                        for r_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                            md_lines.append("| " + " | ".join(cells) + " |")
                            if r_idx == 0:
                                md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        md_lines.append("")

                # 3. Extract Speaker / Presenter Notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        md_lines.append(f"\n> **Speaker Notes**: {notes_text}")

                md_lines.append("\n---\n")

            md_lines.append(f"--- END OF POWERPOINT PRESENTATION ({filename}) ---")
            combined_text = "\n".join(md_lines)

            return ProcessedContent(
                raw_bytes=file_bytes,
                mime_type="text/markdown",
                text_content=combined_text,
                source_filename=filename,
                is_text=True,
                converted_to_pdf=False
            )

        except Exception as e:
            return ProcessedContent(
                raw_bytes=file_bytes,
                mime_type="text/plain",
                text_content=f"PowerPoint presentation parse fallback ({filename}): {str(e)}",
                source_filename=filename,
                is_text=True,
                converted_to_pdf=False
            )
